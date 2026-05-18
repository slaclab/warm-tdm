#!/usr/bin/env python3
"""Generate the ECC AI Seminar demo PowerPoint using the existing template's layout."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load template, extract layout, delete existing slides
prs = Presentation('docs/ECC AI Use Best Practices.pptx')
DEFAULT_LAYOUT = prs.slides[6].slide_layout

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Styling constants (from existing slides)
TITLE_COLOR = RGBColor(0x8F, 0x1D, 0x27)
SUBTITLE_COLOR = RGBColor(0xA8, 0x3E, 0x54)
HEADING_COLOR = RGBColor(0x2B, 0x2B, 0x2B)
BODY_COLOR = RGBColor(0x5C, 0x5C, 0x5C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x1B, 0x3A, 0x5C)
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
TERM_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
MONO_FONT = 'Courier New'


def new_slide(prs):
    return prs.slides.add_slide(DEFAULT_LAYOUT)


def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Emu(548640), Emu(310896), Emu(11064300), Emu(713100))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(548640), Emu(1060704), Emu(11091600), Emu(20100))
    divider.fill.solid()
    divider.fill.fore_color.rgb = TITLE_COLOR
    divider.line.fill.background()


def add_subtitle(slide, text):
    txBox = slide.shapes.add_textbox(Emu(548640), Emu(1130000), Emu(11064300), Emu(400000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = SUBTITLE_COLOR


def add_body(slide, lines, top=Emu(1700000), left=Emu(548640), width=Emu(5500000), height=Emu(5000000)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(opts.get('size', 14))
        p.font.bold = opts.get('bold', False)
        color = opts.get('color', BODY_COLOR)
        if color:
            p.font.color.rgb = color
        if opts.get('mono'):
            p.font.name = MONO_FONT
        if opts.get('indent'):
            p.level = opts['indent']
    return tf


def add_code_block(slide, left, top, width, height, code_lines, title=None):
    """Dark background code block with monospace font."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.02  # Reduce corner rounding
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x55)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000)
    tf.margin_top = Emu(100000)
    tf.margin_right = Emu(150000)
    start = 0
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x99)
        p.font.name = MONO_FONT
        start = 1
    for i, line in enumerate(code_lines):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
        p.font.name = MONO_FONT
    return shape


def add_screenshot_placeholder(slide, left, top, width, height, label):
    """Placeholder for a terminal screenshot."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x55)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000)
    tf.margin_top = Emu(100000)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = TERM_GREEN
    p.font.name = MONO_FONT
    return shape


# ============================================================================
# Slide 1: Title
# ============================================================================
slide = new_slide(prs)
add_title(slide, "AI-Assisted Feature Addition To Existing Code")
add_subtitle(slide, "Generative AI Use Case 2  •  Ben Reese")
add_body(slide, [
    ("", {'size': 6}),
    ("Tool: Claude Code (Anthropic CLI agent)", {'size': 16, 'color': HEADING_COLOR}),
    ("", {'size': 8}),
    ("A 30-minute session converting a fixed-point PID controller", {'size': 14, 'color': BODY_COLOR}),
    ("to floating point — across VHDL, Python, and build infrastructure.", {'size': 14, 'color': BODY_COLOR}),
    ("", {'size': 8}),
    ("• ~800 lines of production code generated", {'size': 14, 'color': BODY_COLOR}),
    ("• 15 files created/modified", {'size': 14, 'color': BODY_COLOR}),
    ("• FPGA firmware + Python drivers + build system", {'size': 14, 'color': BODY_COLOR}),
], top=Emu(1600000), left=Emu(548640), width=Emu(11000000))

# ============================================================================
# Slide 2: The Task
# ============================================================================
slide = new_slide(prs)
add_title(slide, "The Task")
add_subtitle(slide, "Converting a Fixed-Point PID Controller to Floating Point")

add_body(slide, [
    ("The system:", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("• FPGA-based detector readout (VHDL + Python control software)", {'size': 13, 'color': BODY_COLOR}),
    ("• PID servo module: ~1000 lines of VHDL, 8 parallel instances", {'size': 13, 'color': BODY_COLOR}),
    ("• Goal: replace fixed-point math with IEEE 754 float32", {'size': 13, 'color': BODY_COLOR}),
    ("    → Better dynamic range, simpler coefficient tuning", {'size': 12, 'color': BODY_COLOR}),
    ("    → Reuse Xilinx FP IP cores already in the design", {'size': 12, 'color': BODY_COLOR}),
    ("", {'size': 8}),
    ("Why this is a good AI-assisted task:", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("• Not greenfield — must integrate into a large existing codebase", {'size': 13, 'color': BODY_COLOR}),
    ("• Crosses firmware (VHDL) and software (Python) boundaries", {'size': 13, 'color': BODY_COLOR}),
    ("• Requires understanding existing design before writing new code", {'size': 13, 'color': BODY_COLOR}),
    ("• Many files touched, but the design work is the hard part", {'size': 13, 'color': BODY_COLOR}),
], top=Emu(1600000), left=Emu(548640), width=Emu(5800000))

# Right: actual code from AdcDsp.vhd
add_code_block(slide, Emu(6600000), Emu(1600000), Emu(4800000), Emu(4500000), [
    "when PID_P_S =>",
    "  v.pidResult := resize(",
    "    r.pidResult +",
    "    (r.pidCoef * r.pidMultiplier),",
    "    v.pidResult);",
    "  v.pidCoef       := iSfixed;",
    "  v.pidMultiplier := r.sumAccum;",
    "  v.state         := PID_I_S;",
    "",
    "when PID_I_S =>",
    "  v.pidResult := resize(",
    "    r.pidResult +",
    "    (r.pidCoef * r.pidMultiplier),",
    "    v.pidResult);",
    "  v.pidCoef       := dSfixed;",
    "  v.pidMultiplier := resize(",
    "    r.lastAccumError - r.accumError,",
    "    v.pidMultiplier);",
    "  v.state         := PID_D_S;",
    "",
    "when PID_D_S =>",
    "  pidResultNext := resize(",
    "    r.pidResult +",
    "    (r.pidCoef * r.pidMultiplier),",
    "    pidResultNext);",
    "  v.pidResult := pidResultNext;",
], title="── AdcDsp.vhd (existing fixed-point) ──")

# ============================================================================
# Slide 3: The Workflow
# ============================================================================
slide = new_slide(prs)
add_title(slide, "The Workflow")
add_subtitle(slide, "Claude Code operates in the terminal alongside your normal dev tools")

add_body(slide, [
    ("Used via VSCode extension (also", {'size': 13, 'color': BODY_COLOR}),
    ("available as CLI and web app)", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ("Started in Plan Mode:", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("• Agent explores code (read-only)", {'size': 13, 'color': BODY_COLOR}),
    ("• Proposes approach, asks questions", {'size': 13, 'color': BODY_COLOR}),
    ("• No edits until plan is approved", {'size': 13, 'color': BODY_COLOR}),
    ("• Ensures alignment before effort", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ("Then implementation:", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("• Writes code across files", {'size': 13, 'color': BODY_COLOR}),
    ("• Responds to feedback/review", {'size': 13, 'color': BODY_COLOR}),
    ("• Documents progress", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ("~30 min interaction → ~1 day of work", {'size': 13, 'bold': True, 'color': HEADING_COLOR}),
], top=Emu(1600000), left=Emu(548640), width=Emu(5000000))

# Right: terminal box - taller, more content
add_code_block(slide, Emu(5600000), Emu(1200000), Emu(5800000), Emu(5300000), [
    "                    ┌ Plan Mode ┐",
    "╭─────────────────────────────────────────╮",
    "│ > Let's focus on the                    │",
    "│   ColumnFpgaBoard325Coordinator10G      │",
    "│   target. In the data path, there is    │",
    "│   a AdcDsp.vhd block that implements    │",
    "│   a fixed point PID calculation. I'd    │",
    "│   like to adapt this module to use      │",
    "│   floating point math for the PID       │",
    "│   calculation. The design already has   │",
    "│   some floating point math happening    │",
    "│   in it, and a few Vivado IP cores are  │",
    "│   utilized for this. Let's make a plan. │",
    "╰─────────────────────────────────────────╯",
    "",
    "I'll explore the codebase to understand the",
    "existing floating point infrastructure and",
    "the target structure before planning.",
    "",
    "  ⟳ Agent(Explore): existing FP infra",
    "  ⟳ Agent(Explore): target and data path",
    "  ⟳ Agent(Explore): PID state and RAM usage",
    "",
    "Found: FpMac (4-cyc FMA), Int2Fp (2-cyc),",
    "8 instances of AdcDsp in DataPath.vhd,",
    "BiquadFilter.vhd demonstrates the pattern.",
    "DSP48: 8% used, BRAM: 59%, Slices: 75%.",
    "",
    "Writing plan to docs/plans/fp-dsp-pid/...",
], title="── claude (VSCode extension) ──")

# ============================================================================
# Slide 4: AGENTS.md
# ============================================================================
slide = new_slide(prs)
add_title(slide, "AGENTS.md — Giving the Agent Context")
add_subtitle(slide, "The single most impactful thing you can add to a repo for AI assistance")

add_body(slide, [
    ("A Markdown file at the repo root:", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("• Architecture and data flow", {'size': 13, 'color': BODY_COLOR}),
    ("• Where to find things (file/directory map)", {'size': 13, 'color': BODY_COLOR}),
    ("• Naming conventions and coding patterns", {'size': 13, 'color': BODY_COLOR}),
    ("• How to start for each task area", {'size': 13, 'color': BODY_COLOR}),
    ("• Build system patterns", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ("Without → agent guesses for many turns", {'size': 13, 'color': BODY_COLOR}),
    ("With → immediate navigation, correct", {'size': 13, 'bold': True, 'color': HEADING_COLOR}),
    ("conventions from the start", {'size': 13, 'bold': True, 'color': HEADING_COLOR}),
], top=Emu(1600000), left=Emu(548640), width=Emu(4500000))

# Right: one large code block with AGENTS.md excerpts
add_code_block(slide, Emu(5100000), Emu(1500000), Emu(6500000), Emu(4800000), [
    "## Project Summary",
    "Warm TDM is a time-division multiplexing detector",
    "readout system for TES bolometric detector arrays.",
    "",
    "### Data Flow (Column Board)",
    "AD9681 ADC → DataPath → AdcDsp (PID + flux-jump)",
    "  → EventBuilder → PGP Stream → Host",
    "       ↕ FastDacDriver (SQ1 feedback)",
    "",
    "## Essential Reading by Task",
    "| Task Area       | Start With               |",
    "|-----------------|--------------------------|",
    "| DSP / data path | DataPath.vhd, AdcDsp.vhd |",
    "|                 | BiquadFilter.vhd          |",
    "| Timing protocol | TimingPkg.vhd, TimingTx   |",
    "| Communication   | PgpEthCore, RingRouter    |",
    "",
    "## Firmware Conventions",
    "- Library: All RTL loaded as -lib warm_tdm",
    "- VHDL standard: 2008",
    "- Generics suffixed _G (TPD_G, SIMULATION_G)",
    "- Constants suffixed _C (AXIL_CLK_FREQ_C)",
    "- Architecture always named `rtl`",
    "",
    "For substantial feature work, keep planning,",
    "progress, and handoff Markdown under",
    "docs/plans/<task-name>/.",
], title="── AGENTS.md (excerpts) ──")

# ============================================================================
# Slide 5: Plan Docs
# ============================================================================
slide = new_slide(prs)
add_title(slide, "AGENTS.md → Plan Docs")
add_subtitle(slide, "Agent automatically created docs/plans/fp-dsp-pid/PLAN.md")

add_body(slide, [
    ("Why plan docs matter:", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("", {'size': 4}),
    ("• Plans survive context resets", {'size': 13, 'color': BODY_COLOR}),
    ("  → agent resumes in new sessions", {'size': 12, 'color': BODY_COLOR}),
    ("", {'size': 4}),
    ("• Engineers review before code", {'size': 13, 'color': BODY_COLOR}),
    ("  → catch issues early", {'size': 12, 'color': BODY_COLOR}),
    ("", {'size': 4}),
    ("• Decisions + rationale captured", {'size': 13, 'color': BODY_COLOR}),
    ("  → not buried in chat logs", {'size': 12, 'color': BODY_COLOR}),
    ("", {'size': 8}),
    ("This example fit in one session.", {'size': 12, 'color': BODY_COLOR}),
    ("Most features won't — plan docs", {'size': 12, 'color': BODY_COLOR}),
    ("become essential for continuity.", {'size': 12, 'bold': True, 'color': HEADING_COLOR}),
], top=Emu(1600000), left=Emu(548640), width=Emu(4500000))

# Right: the generated plan
add_code_block(slide, Emu(5300000), Emu(1500000), Emu(6200000), Emu(4800000), [
    "# Floating-Point PID (AdcDspFp)",
    "",
    "## Scope",
    "New module, port-compatible with AdcDsp,",
    "selectable via USE_FLOAT_PID_G generic.",
    "",
    "## Architecture",
    "| Core   | Operation   | Latency  |",
    "|--------|-------------|----------|",
    "| FpMac  | A*B+C       | 4 cycles |",
    "| Int2Fp | int → float | 2 cycles |",
    "| Fp2Int | float → int | 2 cycles |",
    "",
    "## Key Design Decisions",
    "1. Track unwrapped feedback value as",
    "   primary state (simplifies flux jump)",
    "2. Pass float directly to downstream",
    "   filter (skip redundant conversion)",
    "3. Iterative flux jump loop",
    "   (handles multi-quantum jumps)",
    "",
    "## State Machine (~38 cycles common case)",
    "ACCUMULATE → Int2Fp → P_FMA → I_FMA →",
    "D_DIFF_FMA → D_FMA → SQ1FB_ADD_FMA →",
    "DERIVE_WRAPPED → Fp2Int → FLUX_CHECK →",
    "ANTI_WINDUP → SUM_UPDATE → DATA_STREAM",
], title="── docs/plans/fp-dsp-pid/PLAN.md ──")

# ============================================================================
# Slide 6: Iteration
# ============================================================================
slide = new_slide(prs)
add_title(slide, "Plan → Iteration")
add_subtitle(slide, "The plan evolved through back-and-forth design dialogue")

add_body(slide, [
    ('"New module or modify in place?"', {'size': 13, 'color': HEADING_COLOR}),
    ("  → New file + generic to select", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ('"Can we pass float directly to the next filter stage?"', {'size': 13, 'color': HEADING_COLOR}),
    ("  → Yes — added bypass generic, eliminates 8 redundant conversions", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ('"How does anti-windup work in float domain?"', {'size': 13, 'color': HEADING_COLOR}),
    ("  → Sign-bit comparison, no extra FP operation needed", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 6}),
    ('"Support multi-quantum flux jumps?"', {'size': 13, 'color': HEADING_COLOR}),
    ("  → Iterative integer loop, widened counter from 9 to 16 bits", {'size': 13, 'color': BODY_COLOR}),
    ("", {'size': 10}),
    ("The agent participates in design — but the engineer", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
    ("makes the architectural calls.", {'size': 14, 'bold': True, 'color': HEADING_COLOR}),
], top=Emu(1600000), left=Emu(548640), width=Emu(5500000))

# Right: terminal-style code block showing the Q&A
add_code_block(slide, Emu(6200000), Emu(1200000), Emu(5200000), Emu(5300000), [
    "┌ Question ────────────────────────────┐",
    "│ For the floating point PID, should   │",
    "│ we create a new module or modify     │",
    "│ AdcDsp.vhd in place?                 │",
    "│                                      │",
    "│ ● New module (AdcDspFp.vhd)          │",
    "│ ○ Modify in place with generic       │",
    "│ ○ Modify in place, replace           │",
    "└────────────────────────────────────────┘",
    "",
    "╭──────────────────────────────────────╮",
    "│ > Can we pass the FP PID results     │",
    "│   directly to the biquad filter?     │",
    "╰──────────────────────────────────────╯",
    "",
    "Yes — add INPUT_IS_FLOAT_G generic to",
    "BiquadFilter. Skip its Int2Fp stage.",
    "Eliminates 8 redundant IP instances.",
    "",
    "╭──────────────────────────────────────╮",
    "│ > Is it possible to have a jump of   │",
    "│   more than 1 flux quantum? Can we   │",
    "│   support this now with FP range?    │",
    "╰──────────────────────────────────────╯",
    "",
    "Yes — iterative integer loop. Each iter",
    "costs 1 cycle. Bounded by DAC range /",
    "fluxQuantum ≈ 4-8 iterations max.",
], title="── claude (plan mode Q&A) ──")

# ============================================================================
# Slide 7: Output
# ============================================================================
slide = new_slide(prs)
add_title(slide, "Implementation Output")
add_subtitle(slide, "All created/modified in one ~30 minute session")

add_body(slide, [
    ("Created:", {'size': 13, 'bold': True, 'color': GREEN}),
    ("firmware/common/warm_tdm/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  rtl/AdcDspFp.vhd", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("  ip/Fp2Int/Fp2Int.xci", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("firmware/python/warm_tdm/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  _AdcDspFp.py", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("", {'size': 6}),
    ("Modified:", {'size': 13, 'bold': True, 'color': BLUE}),
    ("firmware/common/warm_tdm/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  rtl/DataPath.vhd", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("  rtl/BiquadFilter.vhd", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("  rtl/WarmTdmPkg.vhd", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("  ruckus.tcl", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("firmware/targets/ColumnFpgaBoard/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  rtl/ColumnFpgaBoard.vhd", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("firmware/targets/...Coordinator10G/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  ruckus.tcl", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("firmware/python/warm_tdm/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  _DataPath.py, _ColumnFpgaBoard.py", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("  _HardwareGroup.py, __init__.py", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
    ("software/python/warm_tdm_api/", {'size': 11, 'bold': True, 'color': HEADING_COLOR, 'mono': True}),
    ("  _ArgParser.py, _GroupRoot.py, _Group.py", {'size': 11, 'color': BODY_COLOR, 'mono': True}),
], top=Emu(1600000), left=Emu(548640), width=Emu(5400000))

# Right: code sample from the generated VHDL
add_code_block(slide, Emu(5800000), Emu(1500000), Emu(5600000), Emu(4800000), [
    "-- AdcDspFp.vhd (generated)",
    "",
    "when PID_P_S =>",
    "  if (fpMacOutValid = '1') then",
    "    v.pidResultFp := fpMacOutData;",
    "    -- I term: result += I * sumAccum",
    "    v.fpMacInValid := '1';",
    "    v.fpMacA := r.sumAccumFp;",
    "    v.fpMacB := r.iCoef;",
    "    v.fpMacC := fpMacOutData;",
    "    v.state  := PID_I_S;",
    "  end if;",
    "",
    "when PID_I_S =>",
    "  if (fpMacOutValid = '1') then",
    "    v.pidResultFp := fpMacOutData;",
    "    -- D error: lastAccum - accumError",
    "    v.fpMacInValid := '1';",
    "    v.fpMacA := r.accumErrorFp;",
    "    v.fpMacB := X\"BF800000\"; -- -1.0",
    "    v.fpMacC := r.lastAccumErrorFp;",
    "    v.state  := PID_D_DIFF_S;",
    "  end if;",
], title="── firmware/common/warm_tdm/rtl/AdcDspFp.vhd ──")

# ============================================================================
# Slide 8: What Worked / Human
# ============================================================================
slide = new_slide(prs)
add_title(slide, "What Worked / What Needed a Human")

add_body(slide, [
    ("Agent handled well:", {'size': 14, 'bold': True, 'color': GREEN}),
    ("• Reading & comprehending large existing codebase", {'size': 13, 'color': BODY_COLOR}),
    ("• Following established patterns (copied style from BiquadFilter)", {'size': 13, 'color': BODY_COLOR}),
    ("• Cross-domain work (VHDL ↔ Python ↔ TCL build system)", {'size': 13, 'color': BODY_COLOR}),
    ("• Responding to design feedback and updating coherently", {'size': 13, 'color': BODY_COLOR}),
    ("• Auto-generating plan documentation", {'size': 13, 'color': BODY_COLOR}),
], top=Emu(1200000), left=Emu(548640), width=Emu(5500000))

add_body(slide, [
    ("Required human judgment:", {'size': 14, 'bold': True, 'color': TITLE_COLOR}),
    ("• Architecture decisions (new module vs. modify in place)", {'size': 13, 'color': BODY_COLOR}),
    ("• Algorithm correctness (anti-windup semantics, thresholds)", {'size': 13, 'color': BODY_COLOR}),
    ("• Catching missed conventions (review found missing patterns)", {'size': 13, 'color': BODY_COLOR}),
    ("• Deciding what to test and how to verify", {'size': 13, 'color': BODY_COLOR}),
], top=Emu(3400000), left=Emu(548640), width=Emu(5500000))

add_body(slide, [
    ("The AI accelerates implementation;", {'size': 16, 'bold': True, 'color': HEADING_COLOR}),
    ("the engineer owns the design.", {'size': 16, 'bold': True, 'color': HEADING_COLOR}),
], top=Emu(5400000), left=Emu(548640), width=Emu(11000000))

# Right: Claude terminal style showing review catches - taller
add_code_block(slide, Emu(6300000), Emu(1200000), Emu(5100000), Emu(5200000), [
    "╭─────────────────────────────────────╮",
    "│ > _AdcDsp.py has P_CoefRaw,         │",
    "│   I_CoefRaw and D_CoefRaw. The new  │",
    "│   FP version only has I_CoefRaw.    │",
    "│   Why?                              │",
    "╰─────────────────────────────────────╯",
    "",
    "No good reason — it was an oversight.",
    "",
    "  ● Add Raw/Link for all three",
    "  ○ Keep as-is (simplified)",
    "",
    "╭─────────────────────────────────────╮",
    "│ > There should be just one non-     │",
    "│   hidden FluxQuantum LinkVariable.  │",
    "│   It should set all raw registers.  │",
    "╰─────────────────────────────────────╯",
    "",
    "Done. FluxQuantum LinkVariable now sets",
    "FluxQuantumIntRaw, FluxQuantumFpRaw,",
    "and InvFluxQuantumFpRaw internally.",
    "",
    "╭─────────────────────────────────────╮",
    "│ > Make sure the software is also    │",
    "│   updated to instantiate the new    │",
    "│   AdcDspFp device. There should be  │",
    "│   a CLI arg to select FP version.   │",
    "╰─────────────────────────────────────╯",
    "",
    "Threaded --floatPid through:",
    "ArgParser → GroupRoot → Group →",
    "HardwareGroup → ColumnFpgaBoard →",
    "DataPath (selects AdcDspFp vs AdcDsp)",
], title="── claude (review & iteration) ──")

# Note about underspecified software - within the left body area
add_body(slide, [
    ("", {'size': 4}),
    ("Note: Software side was underspecified", {'size': 12, 'bold': True, 'color': HEADING_COLOR}),
    ("in the original prompt and plan —", {'size': 12, 'color': BODY_COLOR}),
    ("additional prompting rounds were needed", {'size': 12, 'color': BODY_COLOR}),
    ("to get Python drivers and CLI flag", {'size': 12, 'color': BODY_COLOR}),
    ("integration generated correctly.", {'size': 12, 'color': BODY_COLOR}),
], top=Emu(5300000), left=Emu(548640), width=Emu(5500000))

# ============================================================================
# Slide 9: Tips
# ============================================================================
slide = new_slide(prs)
add_title(slide, "Practical Tips")

tips = [
    ("1.  Write an AGENTS.md", "Biggest ROI for AI-assisted development"),
    ("2.  Use plan mode", "Get alignment on approach before code is generated"),
    ("3.  Iterate on design", "Treat it as a conversation, not a one-shot prompt"),
    ("4.  Review everything", "Fast code still needs the same scrutiny"),
    ("5.  Track plans in the repo", "Enables multi-session work and handoffs"),
    ("6.  Stay in your domain", "The agent is a force multiplier, not a substitute"),
]

add_body(slide, [
    item
    for title, desc in tips
    for item in [
        (title, {'size': 15, 'bold': True, 'color': HEADING_COLOR}),
        ("     " + desc, {'size': 13, 'color': BODY_COLOR}),
        ("", {'size': 8}),
    ]
], top=Emu(1300000), left=Emu(548640), width=Emu(11000000))

# ============================================================================
# Save
# ============================================================================
output_path = 'docs/plans/fp-dsp-pid/ECC_AI_Demo_Use_Case_2.pptx'
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
